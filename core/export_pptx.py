"""Export PowerPoint — 7 slides d'analyse de portefeuille.

Applique le masque de diapo de l'utilisateur tel quel (pas de barre titre,
pas de numéro de slide). Police : Century Gothic. Nom en bas au centre.
"""

import io

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from core.portfolio import Portfolio
from core import metrics
from core.contribution import performance_contribution, risk_contribution

import os as _os
_TEMPLATE = _os.path.join(_os.path.dirname(_os.path.dirname(__file__)), "template.pptx")
_FONT = "Century Gothic"
_AUTHOR = "FERRAH Mancef"

# Theme colors
_DK2 = RGBColor(0x0E, 0x28, 0x41)
_ACCENT1 = RGBColor(0x15, 0x60, 0x82)
_GREY = RGBColor(0x6B, 0x6B, 0x6B)
_LIGHT_GREY = RGBColor(0xA0, 0xA0, 0xA0)

_SW = Inches(13.33)
_SH = Inches(7.50)


def _clear_slide(slide):
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def _author_footer(slide):
    """Nom centré en bas de chaque slide."""
    box = slide.shapes.add_textbox(
        Inches(4.5), Inches(7.0), Inches(4.33), Inches(0.35),
    )
    p = box.text_frame.paragraphs[0]
    p.text = _AUTHOR
    p.font.name = _FONT
    p.font.size = Pt(8)
    p.font.color.rgb = _LIGHT_GREY
    p.alignment = PP_ALIGN.CENTER


def _chart_img(slide, fig, left, top, width, height):
    w_px = max(400, int(width / 914400 * 96))
    h_px = max(300, int(height / 914400 * 96))
    img = fig.to_image(format="png", width=w_px, height=h_px, scale=2)
    slide.shapes.add_picture(io.BytesIO(img), left, top, width, height)


def _metric_box(slide, x, y, label: str, value: str, w=Inches(2.5)):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.name = _FONT
    p1.font.size = Pt(8)
    p1.font.bold = True
    p1.font.color.rgb = _GREY

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = _FONT
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = _DK2


def _metrics_grid(slide, items, x, y, cols=2, col_w=Inches(2.5), row_h=Inches(0.75)):
    for i, (label, value) in enumerate(items):
        r = i // cols
        c = i % cols
        _metric_box(slide, x + c * col_w, y + r * row_h, label, value, col_w)


def _note(slide, x, y, w, text):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = _FONT
    p.font.size = Pt(8)
    p.font.color.rgb = _LIGHT_GREY


# ── Slides ──────────────────────────────────────────────

def _slide_cover(prs, portfolio, risk_free_rate, rf_desc):
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.placeholder_format and shape.placeholder_format.idx == 0:
            p = shape.text_frame.paragraphs[0]
            p.text = "Analyse de portefeuille"
            p.font.name = _FONT
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = _DK2
        elif shape.placeholder_format and shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.paragraphs[0].text = ", ".join(portfolio.tickers)
            tf.paragraphs[0].font.name = _FONT
            tf.paragraphs[0].font.size = Pt(15)
            tf.paragraphs[0].font.color.rgb = _ACCENT1
            p2 = tf.add_paragraph()
            p2.text = f"Rf : {rf_desc}"
            p2.font.name = _FONT
            p2.font.size = Pt(11)
            p2.font.color.rgb = _GREY
    _author_footer(slide)


def _slide_performance(prs, ptf, bench, rf, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    _chart_img(slide, charts_mod.chart_cumulative_performance(ptf, bench),
               Inches(0.5), Inches(1.0), Inches(8.0), Inches(5.0))

    _metrics_grid(slide, [
        ("RENDEMENT CUMULÉ", f"{metrics.cumulative_return(ptf):.2%}"),
        ("CAGR", f"{metrics.annualized_return(ptf):.2%}"),
        ("ALPHA GÉOMÉTRIQUE", f"{metrics.geometric_alpha(ptf, bench):.2%}"),
        ("BÊTA", f"{metrics.beta(ptf, bench):.2f}"),
        ("VOLATILITÉ ANN.", f"{metrics.annualized_volatility(ptf):.2%}"),
        ("SHARPE", f"{metrics.sharpe_ratio(ptf, rf):.2f}"),
    ], x=Inches(8.8), y=Inches(1.2), cols=2, col_w=Inches(2.2))

    _author_footer(slide)


def _slide_drawdown(prs, ptf, bench, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    _chart_img(slide, charts_mod.chart_drawdown(ptf, bench),
               Inches(0.5), Inches(1.0), Inches(8.0), Inches(5.0))

    _metrics_grid(slide, [
        ("MAX DRAWDOWN PTF", f"{metrics.max_drawdown(ptf):.2%}"),
        ("MAX DRAWDOWN BENCH", f"{metrics.max_drawdown(bench):.2%}"),
        ("DURÉE MAX", f"{metrics.max_drawdown_duration(ptf)} jours"),
        ("TRACKING ERROR", f"{metrics.tracking_error(ptf, bench):.2%}"),
    ], x=Inches(8.8), y=Inches(1.2), cols=2, col_w=Inches(2.2))

    _author_footer(slide)


def _slide_var(prs, ptf, bench, k95, k99, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    _chart_img(slide, charts_mod.chart_return_distribution(ptf, bench),
               Inches(0.5), Inches(1.0), Inches(8.0), Inches(5.0))

    _metrics_grid(slide, [
        ("VaR 95% (1J)", f"{metrics.var_95(ptf):.2%}"),
        ("CVaR 95%", f"{metrics.cvar_95(ptf):.2%}"),
        ("VaR 99% (1J)", f"{metrics.var_99(ptf):.2%}"),
        ("CVaR 99%", f"{metrics.cvar_99(ptf):.2%}"),
    ], x=Inches(8.8), y=Inches(1.2), cols=2, col_w=Inches(2.2))

    v95 = "Adéquat" if k95["model_adequate"] else f"Rejeté (p={k95['p_value']:.4f})"
    v99 = "Adéquat" if k99["model_adequate"] else f"Rejeté (p={k99['p_value']:.4f})"
    _note(slide, Inches(8.8), Inches(4.5), Inches(4.2),
          f"Kupiec 95% : {v95}")
    _note(slide, Inches(8.8), Inches(4.8), Inches(4.2),
          f"Kupiec 99% : {v99}")
    _note(slide, Inches(8.8), Inches(5.1), Inches(4.2),
          "Rolling 504j (Bâle III)")

    _author_footer(slide)


def _slide_garch(prs, ptf, bench, garch_ptf, garch_bench, gjr_ptf, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    _chart_img(slide, charts_mod.chart_garch_volatility(garch_ptf, garch_bench),
               Inches(0.5), Inches(1.0), Inches(7.2), Inches(4.2))

    # GARCH params
    mx = Inches(8.0)
    lbl = slide.shapes.add_textbox(mx, Inches(1.0), Inches(4.5), Inches(0.3))
    p = lbl.text_frame.paragraphs[0]
    p.text = "GARCH(1,1)"
    p.font.name = _FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = _ACCENT1

    _metrics_grid(slide, [
        ("α", f"{garch_ptf['alpha']:.4f}"),
        ("β", f"{garch_ptf['beta']:.4f}"),
        ("PERSISTANCE", f"{garch_ptf['persistence']:.4f}"),
        ("VOL J+1", f"{garch_ptf['forecast_vol_1d']:.3%}"),
    ], x=mx, y=Inches(1.4), cols=2, col_w=Inches(2.5))

    # GJR section
    lbl2 = slide.shapes.add_textbox(mx, Inches(3.7), Inches(4.5), Inches(0.3))
    p2 = lbl2.text_frame.paragraphs[0]
    p2.text = "GJR-GARCH(1,1)"
    p2.font.name = _FONT
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x0F, 0x9E, 0xD5)

    _metrics_grid(slide, [
        ("γ (LEVIER)", f"{gjr_ptf['gamma']:.4f}"),
        ("PERSISTANCE", f"{gjr_ptf['persistence']:.4f}"),
        ("VOL J+1", f"{gjr_ptf['forecast_vol_1d']:.3%}"),
        ("VOL J+10", f"{gjr_ptf['forecast_vol_10d']:.3%}"),
    ], x=mx, y=Inches(4.1), cols=2, col_w=Inches(2.5))

    if gjr_ptf["leverage_effect"]:
        _note(slide, mx, Inches(5.7), Inches(4.5),
              f"Effet de levier confirmé (γ = {gjr_ptf['gamma']:.4f} > 0)")

    _author_footer(slide)


def _slide_regimes(prs, ptf, ms_ptf, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    _chart_img(slide, charts_mod.chart_regime_probabilities(ms_ptf, ptf),
               Inches(0.5), Inches(1.0), Inches(8.0), Inches(5.0))

    calm = ms_ptf["regimes"]["calme"]
    stress = ms_ptf["regimes"]["stress"]

    _metrics_grid(slide, [
        ("P(STRESS)", f"{ms_ptf['current_regime_proba']:.1%}"),
        ("VOL CALME", f"{calm['vol_ann']:.1%} ann."),
        ("VOL STRESS", f"{stress['vol_ann']:.1%} ann."),
        ("DURÉE CALME", f"{ms_ptf['expected_duration_calm']:.0f} j"),
        ("DURÉE STRESS", f"{ms_ptf['expected_duration_stress']:.0f} j"),
        ("REND. CALME", f"{calm['mean_ann']:.1%} ann."),
    ], x=Inches(8.8), y=Inches(1.2), cols=2, col_w=Inches(2.2))

    _author_footer(slide)


def _slide_ratios(prs, ptf, bench, rf, rf_series, charts_mod):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    rf_chart = rf_series if rf_series is not None else rf
    _chart_img(slide, charts_mod.chart_rolling_sharpe(ptf, bench, rf_chart),
               Inches(0.5), Inches(1.0), Inches(8.0), Inches(5.0))

    _metrics_grid(slide, [
        ("SHARPE", f"{metrics.sharpe_ratio(ptf, rf):.2f}"),
        ("SORTINO", f"{metrics.sortino_ratio(ptf, rf):.2f}"),
        ("INFORMATION RATIO", f"{metrics.information_ratio(ptf, bench):.2f}"),
        ("CALMAR", f"{metrics.calmar_ratio(ptf):.2f}"),
    ], x=Inches(8.8), y=Inches(1.2), cols=2, col_w=Inches(2.2))

    _note(slide, Inches(8.8), Inches(4.2), Inches(4.2),
          f"Rf actuel = {rf:.2%}")

    _author_footer(slide)


def _slide_contribution(prs, portfolio, charts_mod):
    """Slide 7 : Contribution perf + décomposition risque côte à côte."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    fig1 = charts_mod.chart_performance_contribution(portfolio)
    fig2 = charts_mod.chart_risk_contribution(portfolio)

    _chart_img(slide, fig1, Inches(0.5), Inches(0.8), Inches(6.0), Inches(5.5))
    _chart_img(slide, fig2, Inches(6.8), Inches(0.8), Inches(5.8), Inches(5.5))

    _author_footer(slide)


def _slide_correlation(prs, portfolio, charts_mod):
    """Slide 8 : Matrice de corrélation (pleine page)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _clear_slide(slide)

    fig = charts_mod.chart_correlation_heatmap(portfolio)

    # Centré sur la slide, grande taille pour lisibilité avec beaucoup d'actifs
    _chart_img(slide, fig, Inches(2.5), Inches(0.6), Inches(8.3), Inches(6.2))

    _author_footer(slide)


# ── Public API ──────────────────────────────────────────

def generate_pptx(portfolio: Portfolio, bench_prices: pd.Series,
                  risk_free_rate: float, rf_series, rf_desc: str,
                  garch_ptf: dict, garch_bench: dict,
                  gjr_ptf: dict,
                  kupiec_95: dict, kupiec_99: dict,
                  ms_ptf: dict) -> bytes:
    from ui import charts as charts_mod

    prs = Presentation(_TEMPLATE)

    ptf = portfolio.cumulative_prices
    common = ptf.index.intersection(bench_prices.index)
    ptf = ptf.loc[common]
    bench = bench_prices.loc[common]

    _slide_cover(prs, portfolio, risk_free_rate, rf_desc)
    _slide_performance(prs, ptf, bench, risk_free_rate, charts_mod)
    _slide_drawdown(prs, ptf, bench, charts_mod)
    _slide_var(prs, ptf, bench, kupiec_95, kupiec_99, charts_mod)
    _slide_garch(prs, ptf, bench, garch_ptf, garch_bench, gjr_ptf, charts_mod)
    _slide_regimes(prs, ptf, ms_ptf, charts_mod)
    _slide_ratios(prs, ptf, bench, risk_free_rate, rf_series, charts_mod)
    _slide_contribution(prs, portfolio, charts_mod)
    _slide_correlation(prs, portfolio, charts_mod)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
